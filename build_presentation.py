"""
build_presentation.py — Generates Bluestock_MF_Presentation.pptx
=================================================================
Creates a 12-slide professional PowerPoint presentation for the
Bluestock Mutual Fund Analytics Capstone Project.

Usage:
    python build_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# --- Bluestock Brand Colors ---
DARK_BLUE   = RGBColor(0x1E, 0x3D, 0x59)   # #1e3d59
ORANGE      = RGBColor(0xFF, 0x6E, 0x40)   # #ff6e40
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY  = RGBColor(0xF8, 0xF9, 0xFA)
MID_GRAY    = RGBColor(0x6C, 0x75, 0x7D)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txBox


def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
    else:
        shape.line.fill.background()
    return shape


def add_image_safe(slide, path, left, top, width, height):
    """Add image only if file exists."""
    if os.path.exists(path):
        slide.shapes.add_picture(path, left, top, width, height)


# ========================
# Slide Builders
# ========================

def slide_title(prs):
    """Slide 1: Title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_bg(slide, DARK_BLUE)
    # Orange accent bar
    add_rect(slide, 0, Inches(4.8), SLIDE_W, Inches(0.08), ORANGE)
    add_textbox(slide, "Mutual Fund Analytics", Inches(1), Inches(1.5),
                Inches(11), Inches(1.2), font_size=44, bold=True,
                color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Bluestock Fintech — N100 Financial Intelligence Platform",
                Inches(1), Inches(2.8), Inches(11), Inches(0.6),
                font_size=20, bold=False, color=ORANGE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Capstone Project I  |  Days 1–7  |  June 2026",
                Inches(1), Inches(5.3), Inches(11), Inches(0.5),
                font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def slide_problem(prs):
    """Slide 2: Problem & Objectives"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GRAY)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), DARK_BLUE)
    add_textbox(slide, "Problem Statement & Objectives", Inches(0.4), Inches(0.25),
                Inches(12), Inches(0.7), font_size=28, bold=True,
                color=WHITE, align=PP_ALIGN.LEFT)
    bullets = [
        "India's mutual fund industry has grown to ₹81+ Lakh Crore AUM with 26 Cr+ folios.",
        "Challenge: No unified analytics view across fund performance, investor demographics,",
        "           SIP trends, and risk metrics.",
        "",
        "Project Objectives:",
        "  1. Build a complete ETL pipeline ingesting 10 datasets into a structured SQLite DB.",
        "  2. Perform EDA and Performance Analytics across 40 fund schemes.",
        "  3. Compute VaR, CVaR, Rolling Sharpe, Sector HHI, and a Fund Scorecard.",
        "  4. Deliver a live interactive Streamlit Dashboard + PDF Report.",
    ]
    add_textbox(slide, "\n".join(bullets), Inches(0.6), Inches(1.5),
                Inches(12.2), Inches(5.5), font_size=15, color=DARK_BLUE)


def slide_data_sources(prs):
    """Slide 3: Data Sources"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GRAY)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), DARK_BLUE)
    add_textbox(slide, "Data Sources & Datasets", Inches(0.4), Inches(0.25),
                Inches(12), Inches(0.7), font_size=28, bold=True, color=WHITE)
    datasets = [
        ("fund_master.csv",            "40 schemes with scheme metadata, expense ratios, risk grades"),
        ("nav_history.csv",            "Daily NAV data for all schemes (2020–2026)"),
        ("investor_transactions.csv",  "32,778 investor transactions (SIP/Lumpsum/Redemption)"),
        ("scheme_performance.csv",     "Risk metrics: Sharpe, Alpha, Beta, Max Drawdown"),
        ("portfolio_holdings.csv",     "Stock-level holdings with sector & weight for equity funds"),
        ("aum_by_fund_house.csv",      "Monthly AUM snapshots by fund house (2022–2025)"),
        ("benchmark_indices.csv",      "NIFTY50, NIFTY100, NIFTY_MIDCAP150, BSE_SMALLCAP"),
        ("category_inflows.csv",       "Monthly net inflows by fund category"),
        ("industry_folio_count.csv",   "Total industry folio count over time"),
        ("monthly_sip_inflows.csv",    "Aggregate monthly SIP industry data"),
    ]
    y = Inches(1.4)
    for i, (name, desc) in enumerate(datasets):
        color = DARK_BLUE if i % 2 == 0 else RGBColor(0x2A, 0x52, 0x75)
        add_rect(slide, Inches(0.4), y, Inches(12.4), Inches(0.48), color)
        add_textbox(slide, f"  {name}", Inches(0.5), y + Pt(4), Inches(3.8), Inches(0.42),
                    font_size=11, bold=True, color=ORANGE)
        add_textbox(slide, desc, Inches(4.2), y + Pt(4), Inches(8.4), Inches(0.42),
                    font_size=11, color=WHITE)
        y += Inches(0.5)


def slide_architecture(prs):
    """Slide 4: Architecture"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GRAY)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), DARK_BLUE)
    add_textbox(slide, "Architecture & Tech Stack", Inches(0.4), Inches(0.25),
                Inches(12), Inches(0.7), font_size=28, bold=True, color=WHITE)

    layers = [
        ("RAW DATA",    "10 CSV Datasets (data/raw/)",               ORANGE),
        ("ETL",         "data_cleaning.py → SQLAlchemy → SQLite (bluestock_mf.db)", RGBColor(0x17, 0x2A, 0x3A)),
        ("ANALYTICS",   "EDA (Seaborn/Plotly) | Performance (SciPy) | Risk (NumPy)", DARK_BLUE),
        ("DASHBOARD",   "Streamlit + Plotly — 4-page interactive web app",            RGBColor(0x2A, 0x52, 0x75)),
        ("OUTPUTS",     "Notebooks (.ipynb) | CSVs | PNGs | PPTX | Final Report",     MID_GRAY),
    ]
    y = Inches(1.5)
    for label, text, color in layers:
        add_rect(slide, Inches(0.6), y, Inches(2.0), Inches(0.75), color)
        add_textbox(slide, label, Inches(0.6), y + Pt(6), Inches(2.0), Inches(0.65),
                    font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        add_textbox(slide, text, Inches(2.8), y + Pt(10), Inches(9.8), Inches(0.55),
                    font_size=13, color=DARK_BLUE)
        y += Inches(1.0)


def slide_eda1(prs):
    """Slide 5: EDA Highlights Part 1"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GRAY)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), DARK_BLUE)
    add_textbox(slide, "EDA Highlights — NAV & AUM", Inches(0.4), Inches(0.25),
                Inches(12), Inches(0.7), font_size=28, bold=True, color=WHITE)
    add_image_safe(slide, "reports/1_nav_trend_analysis.png",
                   Inches(0.3), Inches(1.3), Inches(6.2), Inches(3.0))
    add_image_safe(slide, "reports/2_aum_growth.png",
                   Inches(6.8), Inches(1.3), Inches(6.2), Inches(3.0))
    add_textbox(slide, "NAV trend shows a strong 2023 bull run. SBI Mutual Fund dominates AUM at ₹12.5L+ Cr.",
                Inches(0.4), Inches(4.5), Inches(12.4), Inches(0.8),
                font_size=14, color=DARK_BLUE, align=PP_ALIGN.CENTER)


def slide_eda2(prs):
    """Slide 6: EDA Highlights Part 2"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GRAY)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), DARK_BLUE)
    add_textbox(slide, "EDA Highlights — Investors & SIP", Inches(0.4), Inches(0.25),
                Inches(12), Inches(0.7), font_size=28, bold=True, color=WHITE)
    add_image_safe(slide, "reports/3_sip_inflow_trend.png",
                   Inches(0.3), Inches(1.3), Inches(6.2), Inches(2.8))
    add_image_safe(slide, "reports/5_investor_demographics.png",
                   Inches(6.8), Inches(1.3), Inches(6.2), Inches(2.8))
    add_textbox(slide, "SIP inflows hit an all-time high of ₹31,002 Cr in Dec 2025. Millennials drive the majority of new SIP accounts.",
                Inches(0.4), Inches(4.4), Inches(12.4), Inches(0.8),
                font_size=14, color=DARK_BLUE, align=PP_ALIGN.CENTER)


def slide_perf1(prs):
    """Slide 7: Performance Metrics Part 1"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GRAY)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), DARK_BLUE)
    add_textbox(slide, "Fund Performance Metrics", Inches(0.4), Inches(0.25),
                Inches(12), Inches(0.7), font_size=28, bold=True, color=WHITE)
    metrics = [
        ("CAGR (3yr)", "Best: Axis Small Cap 22.4% | Category Avg: 14.2%"),
        ("Sharpe Ratio", "Best: Mirae Asset Large Cap 1.085 | Rf = 6.5% (RBI Repo)"),
        ("Sortino Ratio", "Penalises only downside volatility — more robust for asymmetric returns"),
        ("Alpha (Annualised)", "Top alpha generators are predominantly Small/Mid Cap funds vs NIFTY100"),
        ("Beta", "Large Cap funds: β ≈ 0.85–1.05 | Small Cap: β ≈ 1.2–1.5"),
        ("Max Drawdown", "Worst drawdown: Axis Small Cap -51.7% (COVID crash reference)"),
    ]
    y = Inches(1.4)
    for label, text in metrics:
        add_rect(slide, Inches(0.5), y, Inches(2.8), Inches(0.65), DARK_BLUE)
        add_textbox(slide, label, Inches(0.55), y + Pt(6), Inches(2.7), Inches(0.55),
                    font_size=13, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
        add_textbox(slide, text, Inches(3.5), y + Pt(10), Inches(9.3), Inches(0.55),
                    font_size=13, color=DARK_BLUE)
        y += Inches(0.82)


def slide_perf2(prs):
    """Slide 8: Performance Metrics Part 2 — Benchmark Comparison"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GRAY)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), DARK_BLUE)
    add_textbox(slide, "Benchmark Comparison & Scorecard", Inches(0.4), Inches(0.25),
                Inches(12), Inches(0.7), font_size=28, bold=True, color=WHITE)
    add_image_safe(slide, "reports/benchmark_comparison.png",
                   Inches(0.3), Inches(1.3), Inches(8.0), Inches(4.2))
    add_textbox(slide, "Fund Scorecard (0–100)",
                Inches(8.6), Inches(1.4), Inches(4.4), Inches(0.4),
                font_size=14, bold=True, color=DARK_BLUE)
    scoring = "30% — 3yr CAGR Rank\n25% — Sharpe Rank\n20% — Alpha Rank\n15% — Expense Ratio (inv.)\n10% — Max Drawdown (inv.)"
    add_textbox(slide, scoring, Inches(8.6), Inches(2.0), Inches(4.4), Inches(2.2),
                font_size=13, color=DARK_BLUE)


def slide_dashboard1(prs):
    """Slide 9: Dashboard Screenshots Part 1"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GRAY)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), DARK_BLUE)
    add_textbox(slide, "Streamlit Dashboard — Industry & Performance Pages",
                Inches(0.4), Inches(0.25), Inches(12), Inches(0.7),
                font_size=26, bold=True, color=WHITE)
    add_textbox(slide, "Page 1: Industry Overview\n• KPI Cards: AUM, SIP, Folios, Schemes\n• AUM Trend Line + AMC Bar Chart",
                Inches(0.5), Inches(1.4), Inches(6.0), Inches(2.0),
                font_size=13, color=DARK_BLUE)
    add_textbox(slide, "Page 2: Fund Performance\n• Risk vs Return Scatter (Bubble = Score)\n• Sortable Fund Scorecard Table\n• NAV Trend with Fund Selector",
                Inches(7.0), Inches(1.4), Inches(6.0), Inches(2.0),
                font_size=13, color=DARK_BLUE)
    add_textbox(slide, "Dashboard is live at: python -m streamlit run dashboard/app.py",
                Inches(0.5), Inches(5.5), Inches(12.0), Inches(0.6),
                font_size=13, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)


def slide_dashboard2(prs):
    """Slide 10: Dashboard Screenshots Part 2"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GRAY)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), DARK_BLUE)
    add_textbox(slide, "Streamlit Dashboard — Investor & Market Trend Pages",
                Inches(0.4), Inches(0.25), Inches(12), Inches(0.7),
                font_size=26, bold=True, color=WHITE)
    add_textbox(slide, "Page 3: Investor Analytics\n• Transaction Amounts by State\n• SIP/Lumpsum/Redemption Split Donut\n• Age Group vs Average SIP Amount",
                Inches(0.5), Inches(1.4), Inches(6.0), Inches(2.5),
                font_size=13, color=DARK_BLUE)
    add_textbox(slide, "Page 4: SIP & Market Trends\n• Dual-Axis: SIP Inflows vs NIFTY 50\n• Category Inflow Heatmap\n• Top 5 Categories by Net Inflow",
                Inches(7.0), Inches(1.4), Inches(6.0), Inches(2.5),
                font_size=13, color=DARK_BLUE)
    add_image_safe(slide, "reports/4_category_inflow_heatmap.png",
                   Inches(0.5), Inches(4.0), Inches(12.0), Inches(2.8))


def slide_findings(prs):
    """Slide 11: Key Findings & Recommendations"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, LIGHT_GRAY)
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.2), DARK_BLUE)
    add_textbox(slide, "Key Findings & Recommendations", Inches(0.4), Inches(0.25),
                Inches(12), Inches(0.7), font_size=28, bold=True, color=WHITE)
    findings = [
        "1. SBI Mutual Fund commands 15%+ of total AUM — concentration risk for the industry.",
        "2. SIP culture is booming — monthly SIP inflows grew 3× between Jan 2022 and Dec 2025.",
        "3. Mirae Asset Large Cap delivers the best risk-adjusted returns (Sharpe: 1.085).",
        "4. ~18% of SIP investors show irregular payment gaps (> 35 days) — at-risk of lapsing.",
        "5. Small Cap funds carry the highest VaR, up to -1.9% daily at the 95% confidence level.",
        "6. T30 cities still dominate 68% of SIP inflows; B30 penetration needs targeted campaigns.",
        "Recommendation: Moderate risk investors should prioritise Mirae Asset Large Cap & Kotak Flexicap.",
        "Recommendation: Investors in High-VaR funds should ensure a 3yr+ investment horizon.",
    ]
    y = Inches(1.4)
    for item in findings:
        color = ORANGE if "Recommendation" in item else DARK_BLUE
        add_textbox(slide, item, Inches(0.6), y, Inches(12.2), Inches(0.55),
                    font_size=13, bold=("Recommendation" in item), color=color)
        y += Inches(0.62)


def slide_thankyou(prs):
    """Slide 12: Thank You"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, DARK_BLUE)
    add_rect(slide, 0, Inches(3.2), SLIDE_W, Inches(0.1), ORANGE)
    add_textbox(slide, "Thank You", Inches(1), Inches(1.5), Inches(11), Inches(1.4),
                font_size=56, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Bluestock Fintech — Capstone Project I: Mutual Fund Analytics",
                Inches(1), Inches(3.5), Inches(11), Inches(0.6),
                font_size=18, color=ORANGE, align=PP_ALIGN.CENTER)
    add_textbox(slide, "GitHub: github.com/binod01nep/MutualFundAnalytics",
                Inches(1), Inches(4.3), Inches(11), Inches(0.5),
                font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    add_textbox(slide, "Dashboard: python -m streamlit run dashboard/app.py",
                Inches(1), Inches(5.0), Inches(11), Inches(0.5),
                font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_title(prs)
    slide_problem(prs)
    slide_data_sources(prs)
    slide_architecture(prs)
    slide_eda1(prs)
    slide_eda2(prs)
    slide_perf1(prs)
    slide_perf2(prs)
    slide_dashboard1(prs)
    slide_dashboard2(prs)
    slide_findings(prs)
    slide_thankyou(prs)

    out_path = "reports/Bluestock_MF_Presentation.pptx"
    prs.save(out_path)
    print(f"Presentation saved to {out_path}  ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
